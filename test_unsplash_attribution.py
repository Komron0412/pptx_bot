"""
Test script to generate a presentation with Unsplash images
and demonstrate the attribution links
"""
import os
from dotenv import load_dotenv
from image_service import ImageService
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Load environment variables
load_dotenv()

def create_test_presentation():
    """Create a test presentation with Unsplash images to show attribution"""
    
    # Initialize image service with Unsplash key
    unsplash_key = os.getenv('UNSPLASH_ACCESS_KEY')
    if not unsplash_key:
        print("ERROR: UNSPLASH_ACCESS_KEY not found in .env file")
        return
    
    image_service = ImageService(unsplash_key=unsplash_key)
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Test queries
    queries = ["nature", "technology", "business"]
    
    for query in queries:
        print(f"\nFetching image for: {query}")
        result = image_service.fetch_image(query)
        
        if result and result.get('credit'):
            credit = result['credit']
            print(f"✓ Image fetched successfully")
            print(f"  Credit text: {credit['text']}")
            print(f"  Photographer link: {credit['link']}")
            print(f"  Unsplash link: {credit.get('app_link', 'N/A')}")
            
            # Create slide
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            
            # Add image
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(9)
            height = Inches(5.5)
            slide.shapes.add_picture(result['path'], left, top, width=width, height=height)
            
            # Add attribution text box
            left = Inches(0.5)
            top = Inches(6.2)
            width = Inches(9)
            height = Inches(1)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            
            # Add the attribution text with hyperlink
            run = p.add_run()
            run.text = credit['text']
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(100, 100, 100)
            
            # Add hyperlink to the photographer's profile
            run.hyperlink.address = credit['link']
            
            print(f"  ✓ Slide created with clickable attribution link")
            
            # Trigger download endpoint for Unsplash compliance
            if result.get('download_url'):
                image_service.trigger_download(result['download_url'])
        else:
            print(f"✗ Failed to fetch image for: {query}")
    
    # Save presentation
    output_path = "generated_presentations/unsplash_attribution_demo.pptx"
    prs.save(output_path)
    print(f"\n✓ Presentation saved to: {output_path}")
    print(f"\nYou can now open this presentation and:")
    print("1. Hover over the attribution text to see the URL in the browser")
    print("2. Click the link to verify it goes to the photographer's Unsplash profile")
    print("3. Take a screenshot showing the URL with UTM parameters")

if __name__ == "__main__":
    create_test_presentation()
