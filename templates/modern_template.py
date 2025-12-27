"""
Modern PowerPoint templates with creative designs
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from templates.styles import ColorScheme, Typography, Layout


class ModernTemplate:
    """Create modern, visually appealing PowerPoint slides"""
    
    def __init__(self, color_scheme=None):
        """Initialize with a color scheme"""
        self.prs = Presentation()
        self.prs.slide_width = Layout.SLIDE_WIDTH
        self.prs.slide_height = Layout.SLIDE_HEIGHT
        self.colors = color_scheme or ColorScheme.get_random_scheme()
    
    def add_title_slide(self, title, subtitle=""):
        """Create a hero title slide with gradient background"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        
        # Add gradient background
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = 45
        fill.gradient_stops[0].color.rgb = self.colors['primary']
        fill.gradient_stops[1].color.rgb = self.colors['secondary']
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Layout.MARGIN_LEFT,
            Inches(2.5),
            Layout.SLIDE_WIDTH - Layout.MARGIN_LEFT - Layout.MARGIN_RIGHT,
            Inches(2.0)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.alignment = PP_ALIGN.CENTER
        # Reduce size to prevent overflow for long titles
        title_p.font.size = Pt(44) if len(title) > 25 else Typography.HERO_SIZE
        title_p.font.bold = True
        title_p.font.color.rgb = self.colors['text']
        title_p.font.name = Typography.TITLE_FONT
        
        # Add subtitle if provided
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Layout.MARGIN_LEFT,
                Inches(4.8),
                Layout.SLIDE_WIDTH - Layout.MARGIN_LEFT - Layout.MARGIN_RIGHT,
                Inches(1.0)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_p = subtitle_frame.paragraphs[0]
            subtitle_p.text = subtitle
            subtitle_p.alignment = PP_ALIGN.CENTER
            subtitle_p.font.size = Typography.SUBTITLE_SIZE
            subtitle_p.font.color.rgb = self.colors['text']
            subtitle_p.font.name = Typography.BODY_FONT
        
        return slide
    
    def add_content_slide(self, title, bullets, image_path=None, credit=None):
        """Create a content slide with optional image"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        
        # Background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['background']
        
        # Title section with accent bar
        title_left = Layout.MARGIN_LEFT
        title_top = Layout.MARGIN_TOP
        title_width = Layout.SLIDE_WIDTH - Layout.MARGIN_LEFT - Layout.MARGIN_RIGHT
        
        # Accent bar
        accent_bar = slide.shapes.add_shape(
            1,  # Rectangle
            title_left,
            title_top,
            Inches(0.15),
            Inches(0.8)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = self.colors['accent']
        accent_bar.line.fill.background()
        
        # Truncate title if too long to prevent overflow
        display_title = title if len(title) < 55 else title[:52] + "..."
        
        # Title text
        title_box = slide.shapes.add_textbox(
            title_left + Inches(0.4), # Increased padding
            title_top,
            title_width - Inches(0.5),
            Inches(1.0)
        )
        title_frame = title_box.text_frame
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        title_p = title_frame.paragraphs[0]
        title_p.text = display_title
        title_p.font.size = Pt(32) if len(display_title) > 30 else Typography.TITLE_SIZE
        title_p.font.bold = True
        title_p.font.color.rgb = self.colors['text']
        title_p.font.name = Typography.TITLE_FONT
        
        # Content area
        content_top = title_top + Inches(1.2)
        
        if image_path:
            # Split layout: content on left, image on right
            content_width = Inches(4.5)
            image_left = Layout.MARGIN_LEFT + content_width + Inches(0.5)
            image_width = Layout.SLIDE_WIDTH - image_left - Layout.MARGIN_RIGHT
            
            # Add bullets
            text_box = slide.shapes.add_textbox(
                Layout.MARGIN_LEFT,
                content_top,
                content_width,
                Inches(5.0)
            )
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            
            for i, bullet in enumerate(bullets[:6]):  # Limit to 6 bullets
                if i > 0:
                    text_frame.add_paragraph()
                p = text_frame.paragraphs[i]
            # Shift content left
            content_width = Inches(4.0)
            
            # Add image on right
            try:
                slide.shapes.add_picture(
                    image_path,
                    Inches(5.5),
                    Inches(1.8),
                    width=Inches(4.0)
                )
            except Exception as e:
                print(f"Error adding image: {e}")

            # Add bullets
            text_box = slide.shapes.add_textbox(
                Layout.MARGIN_LEFT,
                content_top,
                content_width,
                Inches(5.0)
            )
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            
            for i, bullet in enumerate(bullets[:6]):  # Limit to 6 bullets
                if i > 0:
                    text_frame.add_paragraph()
                p = text_frame.paragraphs[i]
                # Truncate long bullet text
                bullet_text = bullet if len(bullet) <= 100 else bullet[:97] + "..."
                p.text = bullet_text
                p.level = 0
                p.font.size = Typography.BODY_SIZE
                p.font.color.rgb = self.colors['text']
                p.font.name = Typography.BODY_FONT
                p.space_before = Pt(12)
        else:
            # Full width content
            content_width = Inches(8.5)
            text_box = slide.shapes.add_textbox(
                Layout.MARGIN_LEFT + Inches(0.5),
                content_top,
                Layout.SLIDE_WIDTH - Layout.MARGIN_LEFT - Layout.MARGIN_RIGHT - Inches(1.0),
                Inches(5.0)
            )
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            
            for i, bullet in enumerate(bullets[:6]):  # Limit to 6 bullets
                if i > 0:
                    text_frame.add_paragraph()
                p = text_frame.paragraphs[i]
                # Truncate long bullet text
                bullet_text = bullet if len(bullet) <= 100 else bullet[:97] + "..."
                p.text = "• " + bullet_text
                p.font.size = Typography.BODY_SIZE
                p.font.color.rgb = self.colors['text']
                p.font.name = Typography.BODY_FONT
                p.space_before = Pt(12)
        
        return slide
    
    def add_image_slide(self, title, image_path, caption=""):
        """Create a full-bleed image slide with overlay text"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        
        # Add full-size image
        try:
            pic = slide.shapes.add_picture(
                image_path,
                0, 0,
                width=Layout.SLIDE_WIDTH,
                height=Layout.SLIDE_HEIGHT
            )
            # Send to back
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element)
        except Exception as e:
            print(f"Could not add image: {e}")
            # Fallback to solid background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = self.colors['primary']
        
        # Add semi-transparent overlay
        overlay = slide.shapes.add_shape(
            1,  # Rectangle
            0,
            Layout.SLIDE_HEIGHT - Inches(2.5),
            Layout.SLIDE_WIDTH,
            Inches(2.5)
        )
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
        overlay.fill.transparency = 0.4
        overlay.line.fill.background()
        
        # Add title on overlay
        title_box = slide.shapes.add_textbox(
            Layout.MARGIN_LEFT,
            Layout.SLIDE_HEIGHT - Inches(2.2),
            Layout.SLIDE_WIDTH - Layout.MARGIN_LEFT - Layout.MARGIN_RIGHT,
            Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.alignment = PP_ALIGN.CENTER
        title_p.font.size = Typography.TITLE_SIZE
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 255, 255)
        title_p.font.name = Typography.TITLE_FONT
        
        if caption:
            caption_box = slide.shapes.add_textbox(
                Layout.MARGIN_LEFT,
                Layout.SLIDE_HEIGHT - Inches(0.8),
                Layout.SLIDE_WIDTH - Layout.MARGIN_LEFT - Layout.MARGIN_RIGHT,
                Inches(0.5)
            )
            caption_frame = caption_box.text_frame
            caption_p = caption_frame.paragraphs[0]
            caption_p.text = caption
            caption_p.alignment = PP_ALIGN.CENTER
            caption_p.font.size = Typography.BODY_SIZE
            caption_p.font.color.rgb = RGBColor(255, 255, 255)
            caption_p.font.name = Typography.BODY_FONT
        
        return slide
    
    def save(self, filename):
        """Save the presentation"""
        self.prs.save(filename)
        return filename
