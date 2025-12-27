"""
Collection of diverse presentation templates with unique designs
Each template provides distinct visual identity and layout styles
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from templates.styles import ColorScheme, Typography, Layout
from templates.modern_template import ModernTemplate
from templates.modern_template import ModernTemplate
import random



def add_varied_image(slide, image_path, credit):
    """Add image in varied positions with geometric shapes"""
    if not image_path:
        return
        
    positions = ['right', 'left', 'top_right']
    pos = random.choice(positions)
    
    try:
        if pos == 'right':
            # Right side image with frame
            pic = slide.shapes.add_picture(
                image_path, Inches(6.5), Inches(5), width=Inches(3)
            )
            # Add frame behind
            frame = slide.shapes.add_shape(
                1, Inches(5.3), Inches(1.3), Inches(4), Inches(3)
            )
            frame.fill.solid()
            frame.fill.fore_color.rgb = RGBColor(200, 200, 200) # Placeholder, ideally needs color scheme access
            frame.line.fill.background()
            # Send frame to back
            slide.shapes._spTree.remove(frame._element)
            slide.shapes._spTree.insert(2, frame._element)
            
        elif pos == 'left':
            # Left side image
            pic = slide.shapes.add_picture(
                image_path, Inches(0.3), Inches(4.5), width=Inches(3)
            )
        elif pos == 'top_right':
            # Smaller top right image
            pic = slide.shapes.add_picture(
                image_path, Inches(6), Inches(2), width=Inches(3.5)
            )
    except Exception as e:
        print(f"Error adding varied image: {e}")


class MinimalistTemplate:
    """Clean, lots of white space, simple lines"""
    
    def __init__(self, color_scheme=None):
        self.prs = Presentation()
        self.prs.slide_width = Layout.SLIDE_WIDTH
        self.prs.slide_height = Layout.SLIDE_HEIGHT
        self.colors = color_scheme or ColorScheme.get_random_scheme()
    
    def add_title_slide(self, title, subtitle=""):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Clean white background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Minimal accent line
        line = slide.shapes.add_shape(
            1, Inches(0.5), Inches(0), Inches(0.2), Layout.SLIDE_HEIGHT
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.colors['primary']
        line.line.fill.background()
        
        # Bottom accent
        shape = slide.shapes.add_shape(
            1, Inches(0.7), Inches(6.5), Inches(2), Inches(0.1)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors['secondary']
        shape.line.fill.background()
        
        # Title - centered below line
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(3.5), Inches(8), Inches(2)
        )
        title_box.text_frame.word_wrap = True
        title_p = title_box.text_frame.paragraphs[0]
        title_p.text = title
        title_p.alignment = PP_ALIGN.CENTER
        title_p.font.size = Pt(40)
        title_p.font.bold = False
        title_p.font.color.rgb = self.colors['primary']
        
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(2), Inches(5.5), Inches(6), Inches(0.8)
            )
            subtitle_p = subtitle_box.text_frame.paragraphs[0]
            subtitle_p.text = subtitle
            subtitle_p.alignment = PP_ALIGN.CENTER
            subtitle_p.font.size = Pt(20)
            subtitle_p.font.color.rgb = self.colors['secondary']
        
        return slide
    
    def add_content_slide(self, title, bullets, image_path=None, credit=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Simple title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.8), Inches(8), Inches(0.7)
        )
        title_box.text_frame.word_wrap = True
        title_p = title_box.text_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(30)
        title_p.font.color.rgb = self.colors['primary']
        
        # Content
        content_top = Inches(2)
        if image_path:
            text_box = slide.shapes.add_textbox(
                Inches(1), content_top, Inches(4), Inches(5)
            )
            try:
                slide.shapes.add_picture(
                    image_path, Inches(5.5), content_top, width=Inches(3.5)
                )
            except: pass
        else:
            text_box = slide.shapes.add_textbox(
                Inches(1.5), content_top, Inches(7), Inches(5)
            )
        
        text_box.text_frame.word_wrap = True
        
        # Limit bullets to 6
        for i, bullet in enumerate(bullets[:6]):
            if i > 0:
                text_box.text_frame.add_paragraph()
            p = text_box.text_frame.paragraphs[i]
            # Truncate long bullet text
            bullet_text = bullet if len(bullet) <= 100 else bullet[:97] + "..."
            p.text = "• " + bullet_text
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(50, 50, 50) # Enforce dark text for white background
            p.space_before = Pt(12)
        
        return slide
    
    def save(self, filename):
        self.prs.save(filename)
        return filename


class BoldModernTemplate:
    """Large text, strong colors, asymmetric layouts"""
    
    def __init__(self, color_scheme=None):
        self.prs = Presentation()
        self.prs.slide_width = Layout.SLIDE_WIDTH
        self.prs.slide_height = Layout.SLIDE_HEIGHT
        self.colors = color_scheme or ColorScheme.get_random_scheme()
    
    def add_title_slide(self, title, subtitle=""):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Bold diagonal background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['primary']
        
        # Large geometric shape
        shape = slide.shapes.add_shape(
            1, Inches(-1), Inches(4), Inches(7), Inches(5)
        )
        shape.rotation = 15
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors['secondary']
        shape.fill.transparency = 0.3
        shape.line.fill.background()
        
        # Bold title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(3)
        )
        title_box.text_frame.word_wrap = True
        title_p = title_box.text_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(44)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 255, 255)
        
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1.2), Inches(5.2), Inches(7), Inches(1)
            )
            subtitle_p = subtitle_box.text_frame.paragraphs[0]
            subtitle_p.text = subtitle
            subtitle_p.font.size = Pt(24)
            subtitle_p.font.color.rgb = self.colors['accent']
        
        return slide
    
    def add_content_slide(self, title, bullets, image_path=None, credit=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['primary']  # Colorful background
        
        # White text for contrast on dark/colorful background
        text_color = RGBColor(255, 255, 255)
        accent_color = RGBColor(255, 255, 255) # White accent for contrast
        bar = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(0.3), Layout.SLIDE_HEIGHT
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.colors['accent']
        bar.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(0.5), Inches(9), Inches(1)
        )
        title_box.text_frame.word_wrap = True
        title_p = title_box.text_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(32)
        title_p.font.bold = True
        title_p.font.color.rgb = accent_color
        
        # Content
        if image_path:
            text_box = slide.shapes.add_textbox(
                Inches(0.6), Inches(1.8), Inches(4.5), Inches(5)
            )
            try:
                slide.shapes.add_picture(
                    image_path, Inches(5.5), Inches(1.8), width=Inches(4)
                )
            except: pass
        else:
            text_box = slide.shapes.add_textbox(
                Inches(0.6), Inches(1.8), Inches(8.5), Inches(5)
            )
        
        text_box.text_frame.word_wrap = True
        
        for i, bullet in enumerate(bullets[:6]):
            if i > 0:
                text_box.text_frame.add_paragraph()
            p = text_box.text_frame.paragraphs[i]
            bullet_text = bullet if len(bullet) <= 100 else bullet[:97] + "..."
            p.text = "▸ " + bullet_text
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = text_color
            p.space_before = Pt(14)
        
        return slide
    
    def save(self, filename):
        self.prs.save(filename)
        return filename


class CorporateTemplate:
    """Professional, structured, business-focused"""
    
    def __init__(self, color_scheme=None):
        self.prs = Presentation()
        self.prs.slide_width = Layout.SLIDE_WIDTH
        self.prs.slide_height = Layout.SLIDE_HEIGHT
        self.colors = color_scheme or ColorScheme.get_random_scheme()
    
    def add_title_slide(self, title, subtitle=""):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Two-tone background
        top_rect = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Layout.SLIDE_WIDTH, Inches(4)
        )
        top_rect.fill.solid()
        top_rect.fill.fore_color.rgb = self.colors['primary']
        top_rect.line.fill.background()
        
        bottom_rect = slide.shapes.add_shape(
            1, Inches(0), Inches(4), Layout.SLIDE_WIDTH, Inches(3.5)
        )
        bottom_rect.fill.solid()
        bottom_rect.fill.fore_color.rgb = RGBColor(245, 245, 245)
        bottom_rect.line.fill.background()
        
        # Title on top section
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.5), Inches(8), Inches(2)
        )
        title_box.text_frame.word_wrap = True
        title_p = title_box.text_frame.paragraphs[0]
        title_p.text = title
        title_p.alignment = PP_ALIGN.CENTER
        title_p.font.size = Pt(40)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 255, 255)
        
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(4.5), Inches(7), Inches(1.5)
            )
            subtitle_p = subtitle_box.text_frame.paragraphs[0]
            subtitle_p.text = subtitle
            subtitle_p.alignment = PP_ALIGN.CENTER
            subtitle_p.font.size = Pt(24)
            subtitle_p.font.color.rgb = self.colors['primary']
        
        return slide
    
    def add_content_slide(self, title, bullets, image_path=None, credit=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(250, 250, 250)
        
        # Sidebar
        sidebar = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(0.5), Layout.SLIDE_HEIGHT
        )
        sidebar.fill.solid()
        sidebar.fill.fore_color.rgb = self.colors['secondary']
        sidebar.line.fill.background()
        
        # Header bar
        header = slide.shapes.add_shape(
            1, Inches(0.5), Inches(0), Layout.SLIDE_WIDTH - Inches(0.5), Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.colors['primary']
        header.line.fill.background()
        
        # Title in header
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
        )
        title_box.text_frame.word_wrap = True
        title_p = title_box.text_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(28)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Content area
        if image_path:
            text_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.8), Inches(4.5), Inches(5)
            )
            try:
                slide.shapes.add_picture(
                    image_path, Inches(5.5), Inches(2), width=Inches(4)
                )
            except: pass
        else:
            text_box = slide.shapes.add_textbox(
                Inches(1), Inches(2), Inches(8), Inches(5)
            )
        
        text_box.text_frame.word_wrap = True
        
        for i, bullet in enumerate(bullets[:6]):
            if i > 0:
                text_box.text_frame.add_paragraph()
            p = text_box.text_frame.paragraphs[i]
            bullet_text = bullet if len(bullet) <= 100 else bullet[:97] + "..."
            p.text = "■ " + bullet_text
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(50, 50, 50) # Enforce dark text for white background
            p.space_before = Pt(12)
        
        return slide
    
    def save(self, filename):
        self.prs.save(filename)
        return filename


class CreativeTemplate:
    """Artistic shapes, overlapping elements, dynamic"""
    
    def __init__(self, color_scheme=None):
        self.prs = Presentation()
        self.prs.slide_width = Layout.SLIDE_WIDTH
        self.prs.slide_height = Layout.SLIDE_HEIGHT
        self.colors = color_scheme or ColorScheme.get_random_scheme()
    
    def add_title_slide(self, title, subtitle=""):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = 135
        fill.gradient_stops[0].color.rgb = self.colors['background']
        fill.gradient_stops[1].color.rgb = self.colors['primary']
        
        # Creative circles
        for i in range(3):
            circle = slide.shapes.add_shape(
                3, Inches(6 + i), Inches(0.5 + i * 1.5), Inches(2.5), Inches(2.5)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = self.colors['accent']
            circle.fill.transparency = 0.6
            circle.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(2.5), Inches(7), Inches(2.5)
        )
        title_box.text_frame.word_wrap = True
        title_p = title_box.text_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(42)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 255, 255)
        
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(5.2), Inches(6), Inches(1)
            )
            subtitle_p = subtitle_box.text_frame.paragraphs[0]
            subtitle_p.text = subtitle
            subtitle_p.font.size = Pt(22)
            subtitle_p.font.color.rgb = self.colors['accent']
        
        return slide
    
    def add_content_slide(self, title, bullets, image_path=None, credit=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['background']
        
        # Creative shape accent
        shape = slide.shapes.add_shape(
            3, Inches(8), Inches(-0.5), Inches(3), Inches(3)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors['accent']
        shape.fill.transparency = 0.7
        shape.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(0.5), Inches(8), Inches(0.9)
        )
        title_box.text_frame.word_wrap = True
        title_p = title_box.text_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(30)
        title_p.font.bold = True
        title_p.font.color.rgb = self.colors['accent']
        
        # Content
        if image_path:
            text_box = slide.shapes.add_textbox(
                Inches(0.6), Inches(1.7), Inches(4.5), Inches(5)
            )
            try:
                slide.shapes.add_picture(
                    image_path, Inches(5.5), Inches(1.7), width=Inches(4)
                )
            except: pass
        else:
            text_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(1.7), Inches(8), Inches(5)
            )
        
        text_box.text_frame.word_wrap = True
        
        for i, bullet in enumerate(bullets[:6]):
            if i > 0:
                text_box.text_frame.add_paragraph()
            p = text_box.text_frame.paragraphs[i]
            bullet_text = bullet if len(bullet) <= 100 else bullet[:97] + "..."
            p.text = "→ " + bullet_text
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['text']
            p.space_before = Pt(12)
        
        return slide
    
    def save(self, filename):
        self.prs.save(filename)
        return filename


class ElegantTemplate:
    """Refined, subtle gradients, serif-inspired"""
    
    def __init__(self, color_scheme=None):
        self.prs = Presentation()
        self.prs.slide_width = Layout.SLIDE_WIDTH
        self.prs.slide_height = Layout.SLIDE_HEIGHT
        self.colors = color_scheme or ColorScheme.get_random_scheme()
    
    def add_title_slide(self, title, subtitle=""):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = 90
        fill.gradient_stops[0].color.rgb = RGBColor(255, 255, 255)
        fill.gradient_stops[1].color.rgb = RGBColor(240, 240, 245)
        
        # Elegant border
        border = slide.shapes.add_shape(
            1, Inches(0.8), Inches(2), Inches(8.4), Inches(3.5)
        )
        border.fill.background()
        border.line.color.rgb = self.colors['primary']
        border.line.width = Pt(2)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1.2), Inches(2.5), Inches(7.6), Inches(2)
        )
        title_box.text_frame.word_wrap = True
        title_p = title_box.text_frame.paragraphs[0]
        title_p.text = title
        title_p.alignment = PP_ALIGN.CENTER
        title_p.font.size = Pt(40)
        title_p.font.color.rgb = self.colors['primary']
        
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(4.5), Inches(7), Inches(0.8)
            )
            subtitle_p = subtitle_box.text_frame.paragraphs[0]
            subtitle_p.text = subtitle
            subtitle_p.alignment = PP_ALIGN.CENTER
            subtitle_p.font.size = Pt(20)
            subtitle_p.font.color.rgb = self.colors['secondary']
            subtitle_p.font.italic = True
        
        return slide
    
    def add_content_slide(self, title, bullets, image_path=None, credit=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        background = slide.background
        fill = background.fill
        # Make background slightly tinted for contrast
        fill.solid()
        fill.fore_color.rgb = RGBColor(248, 248, 250)
        
        # Minimal elegant line
        line = slide.shapes.add_shape(
            1, Inches(1), Inches(1.2), Inches(8), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.colors['primary']
        line.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.4), Inches(8), Inches(0.8)
        )
        title_box.text_frame.word_wrap = True
        title_p = title_box.text_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(32)
        title_p.font.italic = True
        title_p.font.name = Typography.TITLE_FONT
        title_p.font.color.rgb = self.colors['primary']
        
        # Content
        if image_path:
            text_box = slide.shapes.add_textbox(
                Inches(1), Inches(1.5), Inches(4.5), Inches(5)
            )
            try:
                slide.shapes.add_picture(
                    image_path, Inches(6), Inches(1.5), width=Inches(3)
                )
            except: pass
        else:
            text_box = slide.shapes.add_textbox(
                Inches(1), Inches(2), Inches(8), Inches(5)
            )
        
        text_box.text_frame.word_wrap = True
        
        for i, bullet in enumerate(bullets[:6]):
            if i > 0:
                text_box.text_frame.add_paragraph()
            p = text_box.text_frame.paragraphs[i]
            bullet_text = bullet if len(bullet) <= 100 else bullet[:97] + "..."
            p.text = "◆ " + bullet_text
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(50, 50, 50) # Dark gray text, not schema dependent often too light
            p.space_before = Pt(12)
        
        return slide
    
    def save(self, filename):
        self.prs.save(filename)
        return filename


class GeometricTemplate:
    """Shapes and patterns, bold geometry"""
    
    def __init__(self, color_scheme=None):
        self.prs = Presentation()
        self.prs.slide_width = Layout.SLIDE_WIDTH
        self.prs.slide_height = Layout.SLIDE_HEIGHT
        self.colors = color_scheme or ColorScheme.get_random_scheme()
    
    def add_title_slide(self, title, subtitle=""):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['background']
        
        # Geometric triangles
        triangle1 = slide.shapes.add_shape(
            5, Inches(0), Inches(0), Inches(4), Inches(4)
        )
        triangle1.fill.solid()
        triangle1.fill.fore_color.rgb = self.colors['primary']
        triangle1.line.fill.background()
        
        triangle2 = slide.shapes.add_shape(
            5, Inches(7), Inches(4), Inches(3), Inches(3.5)
        )
        triangle2.fill.solid()
        triangle2.fill.fore_color.rgb = self.colors['accent']
        triangle2.fill.transparency = 0.5
        triangle2.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.5), Inches(7), Inches(2.5)
        )
        title_box.text_frame.word_wrap = True
        title_p = title_box.text_frame.paragraphs[0]
        title_p.text = title
        title_p.alignment = PP_ALIGN.CENTER
        title_p.font.size = Pt(42)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 255, 255)
        
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(2), Inches(5.2), Inches(6), Inches(1)
            )
            subtitle_p = subtitle_box.text_frame.paragraphs[0]
            subtitle_p.text = subtitle
            subtitle_p.alignment = PP_ALIGN.CENTER
            subtitle_p.font.size = Pt(22)
            subtitle_p.font.color.rgb = self.colors['accent']
        
        return slide
    
    def add_content_slide(self, title, bullets, image_path=None, credit=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['background']
        
        # Geometric corner accent
        corner = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(2), Inches(1.5)
        )
        corner.fill.solid()
        corner.fill.fore_color.rgb = self.colors['accent']
        corner.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(9), Inches(0.8)
        )
        title_box.text_frame.word_wrap = True
        title_p = title_box.text_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(30)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Content and Varied Image
        content_width = Inches(8.5)
        content_left = Inches(0.7)
        
        if image_path:
            # Pick a random layout
            layouts = ['right', 'left', 'bottom']
            layout = random.choice(layouts)
            
            try:
                if layout == 'right':
                    # Image Right
                    content_width = Inches(4.5)
                    text_box = slide.shapes.add_textbox(
                        Inches(0.7), Inches(1.7), content_width, Inches(5)
                    )
                    
                    # Geometric shape behind image
                    shape = slide.shapes.add_shape(
                        5, Inches(5.8), Inches(1.5), Inches(3.5), Inches(3.5)
                    )
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = self.colors['primary']
                    shape.line.fill.background()
                    
                    pic = slide.shapes.add_picture(
                        image_path, Inches(6), Inches(1.7), width=Inches(3)
                    )
                    
                elif layout == 'left':
                    # Image Left
                    move_right = Inches(4)
                    content_left = Inches(0.7) + move_right
                    content_width = Inches(4.5)
                    
                    # Geometric shape behind image
                    shape = slide.shapes.add_shape(
                        5, Inches(0.8), Inches(1.5), Inches(3.5), Inches(3.5)
                    )
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = self.colors['primary']
                    shape.line.fill.background()

                    pic = slide.shapes.add_picture(
                        image_path, Inches(1), Inches(1.7), width=Inches(3)
                    )

                    text_box = slide.shapes.add_textbox(
                         content_left, Inches(1.7), content_width, Inches(5)
                    )

                else: # Bottom
                    # Image Bottom
                    text_box = slide.shapes.add_textbox(
                        Inches(0.7), Inches(1.5), Inches(8.5), Inches(2.5)
                    )
                    
                    pic = slide.shapes.add_picture(
                        image_path, Inches(3), Inches(4.2), width=Inches(4)
                    )
                    
            except Exception as e: 
                print(f"Error adding varied image: {e}")
                text_box = slide.shapes.add_textbox(
                    Inches(0.7), Inches(1.7), Inches(8.5), Inches(5.5)
                )
        else:
            text_box = slide.shapes.add_textbox(
                Inches(0.7), Inches(1.7), Inches(8.5), Inches(5.5)
            )
        
        text_box.text_frame.word_wrap = True
        
        for i, bullet in enumerate(bullets[:6]):
            if i > 0:
                text_box.text_frame.add_paragraph()
            p = text_box.text_frame.paragraphs[i]
            bullet_text = bullet if len(bullet) <= 100 else bullet[:97] + "..."
            p.text = "▪ " + bullet_text
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['text']
            p.space_before = Pt(12)
        
        return slide
    
    def save(self, filename):
        self.prs.save(filename)
        return filename


def get_random_template():
    """Select a random template class"""
    templates = [
        ModernTemplate,
        MinimalistTemplate,
        BoldModernTemplate,
        CorporateTemplate,
        CreativeTemplate,
        ElegantTemplate,
        GeometricTemplate,
    ]
    return random.choice(templates)
